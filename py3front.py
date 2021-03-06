
import csv
import scipy as sp
#import matplotlib.pyplot as plt
from sklearn.linear_model import LinearRegression
import numpy as np
from flask import Flask
from flask import jsonify,request,redirect,url_for,send_from_directory,render_template,url_for,send_file,make_response, session
from flaskext.mysql import MySQL
import json
import os
from werkzeug import secure_filename
import pefile
import shutil
from sklearn.linear_model import ElasticNet
from sklearn.linear_model import ElasticNetCV
from sklearn.cross_validation import KFold
#from sklearn.model_selection import KFold
from sklearn.feature_extraction.text import CountVectorizer
import sys
import nltk.stem
import docx2txt
from pptx import Presentation
import requests
from time import sleep
import datetime
#from  datetime import datetime
from konlpy.tag import Twitter
from konlpy.utils import pprint
import sklearn
import nltk
#from nltk.book import *
import math
import re
from sklearn.cluster import KMeans
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics import silhouette_samples, silhouette_score

#-*- coding: utf-8 -*-

mystopword = frozenset([
    "가","가운데","갈","걔","거","건","것","걸","고","곳","곳곳","그","그간","그것","그곳","그녀","그달","그당시","그대","그대신",
   "그동안","그들","그들대로","그때","그런고","그런날","그런데서","그런줄","그럴수록","그로","그무렵","그외","그이","그전","그전날",
   "그쪽","나","내게","내년초","내달","내부",
   "너","너로","너와","너희","너희대로","너희들","네","네번","네째","네탓","넷","넷째","년","년간","년도","녘","놈","누가","누구","누구누구",
   "누군가","다섯째","다음달","다음주","당분간","대다수","덧","데","되","둘","둘째","둥","뒤","뒷받침",
   "듯","등","따름","따위","딴","때","때문","마련","마지막","마찬가지","마리","만원","만큼","맏","무엇","몇","묶음","물론","못","뭣","밑",
   "밖에","백","백만원","뿐","서로","세","세째","수십","스스로","쉰","십","아홉째","어느편","어디","어디론지","어떤때",
   "여러가지","여섯째","역","열째","옆","요즘","우리","우선","움직임","월","위해서","육","으뜸",
   "이곳","이것","이날","이달초","이듬","이듬달","이때","이런저런","이런줄","이번","일곱째","쟤",
   "저것","저곳","저도","저기","저런날","저런줄","저렴","저쪽",
   "저희","적극적","전날","전년","전부","전부문","전체적","제","줄곳","지난해",
   "첫째","첫날","최근","평상시","하나","허나","한마디","한가운데","한가지","한곳","한번",
   "한쪽","한편","할",
   "a", "about", "above", "across", "after", "afterwards", "again", "against",
    "all", "almost", "alone", "along", "already", "also", "although", "always",
    "am", "among", "amongst", "amoungst", "amount", "an", "and", "another",
    "any", "anyhow", "anyone", "anything", "anyway", "anywhere", "are",
    "around", "as", "at", "back", "be", "became", "because", "become",
    "becomes", "becoming", "been", "before", "beforehand", "behind", "being",
    "below", "beside", "besides", "between", "beyond", "bill", "both",
    "bottom", "but", "by", "call", "can", "cannot", "cant", "co", "con",
    "could", "couldnt", "cry", "de", "describe", "detail", "do", "done",
    "down", "due", "during", "each", "eg", "eight", "either", "eleven", "else",
    "elsewhere", "empty", "enough", "etc", "even", "ever", "every", "everyone",
    "everything", "everywhere", "except", "few", "fifteen", "fifty", "fill",
    "find", "fire", "first", "five", "for", "former", "formerly", "forty",
    "found", "four", "from", "front", "full", "further", "get", "give", "go",
    "had", "has", "hasnt", "have", "he", "hence", "her", "here", "hereafter",
    "hereby", "herein", "hereupon", "hers", "herself", "him", "himself", "his",
    "how", "however", "hundred", "i", "ie", "if", "in", "inc", "indeed",
    "interest", "into", "is", "it", "its", "itself", "keep", "last", "latter",
    "latterly", "least", "less", "ltd", "made", "many", "may", "me",
    "meanwhile", "might", "mill", "mine", "more", "moreover", "most", "mostly",
    "move", "much", "must", "my", "myself", "name", "namely", "neither",
    "never", "nevertheless", "next", "nine", "no", "nobody", "none", "noone",
    "nor", "not", "nothing", "now", "nowhere", "of", "off", "often", "on",
    "once", "one", "only", "onto", "or", "other", "others", "otherwise", "our",
    "ours", "ourselves", "out", "over", "own", "part", "per", "perhaps",
    "please", "put", "rather", "re", "same", "see", "seem", "seemed",
    "seeming", "seems", "serious", "several", "she", "should", "show", "side",
    "since", "sincere", "six", "sixty", "so", "some", "somehow", "someone",
    "something", "sometime", "sometimes", "somewhere", "still", "such",
    "system", "take", "ten", "than", "that", "the", "their", "them",
    "themselves", "then", "thence", "there", "thereafter", "thereby",
    "therefore", "therein", "thereupon", "these", "they", "thick", "thin",
    "third", "this", "those", "though", "three", "through", "throughout",
    "thru", "thus", "to", "together", "too", "top", "toward", "towards",
    "twelve", "twenty", "two", "un", "under", "until", "up", "upon", "us",
    "very", "via", "was", "we", "well", "were", "what", "whatever", "when",
    "whence", "whenever", "where", "whereafter", "whereas", "whereby",
    "wherein", "whereupon", "wherever", "whether", "which", "while", "whither",
    "who", "whoever", "whole", "whom", "whose", "why", "will", "with",
    "within", "without", "would", "yet", "you", "your", "yours", "yourself",
    "yourselves"])
TXT_DIRECTORY= '/home/gxicxigouxa/myproject/textcompare'
UPLOAD_FOLDER= '/home/gxicxigouxa/myproject' # path convert!! (when i use different server)
ALLOWED_EXTENSIONS =    set(['txt','pdf','png','jpg','jpeg','gif','doc','exe','docx','hwp','pptx'])
OPENSTACK_IP = "183.103.47.19"

#app = Flask(__name__,template_folder='templates', static_folder='/home/asdiste/myproject/templates')
app = Flask(__name__, static_path='/static')
app.config['UPLOAD_FOLDER']=UPLOAD_FOLDER
app.config['TXT_DIRECTORY']=TXT_DIRECTORY
app.debug=True
mysql= MySQL()

app.config['MYSQL_DATABASE_USER']='root'
app.config['MYSQL_DATABASE_PASSWORD']='root'
app.config['MYSQL_DATABASE_DB']='jhj'
app.secret_key = "OPENSTACK_SECRET_KEY"
mysql.init_app(app)

english_stemmer=nltk.stem.SnowballStemmer('english')
class StemmedCountVectorizer(CountVectorizer):
   def build_analyzer(self):
      analyzer = super(StemmedCountVectorizer,self).build_analyzer()
      return lambda doc: (english_stemmer.stem(w) for w in analyzer(doc))

def get_file_list(path):
	filelist=[]

def allowed_file(filename):
   return '.' in filename and \
      filename.rsplit('.',1)[1] in ALLOWED_EXTENSIONS

def dist_norm(v1,v2):
   v1_normalized = v1/sp.linalg.norm(v1.toarray())
   v2_normalized = v2/sp.linalg.norm(v2.toarray())
   delta = v1_normalized -v2_normalized
   return sp.linalg.norm(delta.toarray())

def EnLrswap (enprod, lrprod):
	i = 0
	for i in range(len(enprod)):
		if(enprod[i] < 0):
			enprod[i] = lrprod[i]

def num_only_file(pwd):
   numOfFile=0
   filenames = os.listdir(pwd)
   for filename in filenames:
      if (os.path.isfile(pwd+'/'+filename)):
         numOfFile+=1
         
   return numOfFile  

@app.route('/login')
def showloginpage():
	#######get the admin token
	headers = {"content-type":"application/json"}
	data= {"auth":{"tenantName":"admin","passwordCredentials":{"username":"admin","password":"openstack"}}}
	print(json.dumps(data))
	admin_token_response = requests.post('http://' + OPENSTACK_IP + ':5000/v2.0/tokens',data=json.dumps(data),headers=headers)
	print(admin_token_response)
	json_dict= json.loads(admin_token_response.text)
	admin_token=json_dict['access']['token']['id']
	print("session adminToken: " + admin_token)
	session["adminToken"] = admin_token
	return render_template('oldlogin.html',err_code ="no error")

@app.route('/loginchk', methods=['POST','GET'])
def chklogin():
	if request.method == "POST":
		userId= request.form['id']
		userPwd= request.form['password']
		current_time = datetime.datetime.now()
		print(current_time)

		print (userId)
		conn =mysql.connect()
		cursor =conn.cursor()
		query = "select * from userinfotable where id ='"+userId +"';"
		cursor.execute(query)
		conn.commit()

		   
		result =[]
		columns= tuple( [d[0] for d in cursor.description] )

		   
		for row in cursor:
		   result.append(dict(zip(columns, row)))
   
		print (str(result))

		if (str(result)=='[]'): # When there is no user id in db
			return render_template('oldlogin.html',login_err_code ="id incorrect", sign_up_err_code = "none")
		else:
			convertedDatetime = datetime.datetime.strptime(result[0]["expiretime"], '%Y-%m-%d %H:%M:%S')
			if result[0]["id"] == "admin":
				print("rendering manager monitoring page...")
				return redirect("/monitoring")
			elif(result[0]["pwd"]==userPwd and (convertedDatetime > current_time)):#pwd is correct
				#토큰 가져오기
				print("after compare")		
				token_url = 'http://' + OPENSTACK_IP + ':5000/v2.0/tokens'
				data = {"auth":{"tenantName":userId,"passwordCredentials":{"username":userId,"password":userPwd}}}
				headers = {'content-type':'application/json'}
				response = requests.post(url=token_url,data=json.dumps(data),headers=headers)
				json_data = json.loads(response.text)

				token=json_data["access"]["token"]["id"]
				
				print("token: " + token)
				session["token"] = token
				session["userId"] = userId
				print("after session!")
				#토큰 계속 요청하면 문제 생길 수 있으므로 임의의 토큰, 스토리지 목록을 만들어 보내자.
				'''
				session["token"] = "0123123myuserrandomtoken3213210"
				session["userId"] = userId
				'''
				return redirect("/storage")
			elif (result[0]["pwd"]==userPwd and convertedDatetime < current_time):
				return render_template('oldlogin.html', login_err_code="expired", sign_up_err_code = "none")
			else:#pwd is incorrect
				return render_template('oldlogin.html', login_err_code="pwd incorrect", sign_up_err_code = "none")


@app.route('/signupchk', methods=['POST','GET'])
def chksignup():
	if request.method == "POST":
		signUpId= request.form['user-id']
		signUpPwd= request.form['user-password']
		signUpBirthday = request.form['user-birthday']
		signUpEmail = request.form['user-email']
		paymentDay = 30
		#expiretime = datetime.datetime.now()
		expiretime = datetime.datetime.now() + datetime.timedelta(days=paymentDay)
		print (signUpId)
		rating = 0
		totalamount = 0
		enrollmentnumber = 0
		yearaverageamount = 0
		conn =mysql.connect()
		cursor =conn.cursor()
		query = "select * from userinfotable where id ='"+signUpId +"';"
		cursor.execute(query)
		conn.commit()

		   
		result =[]
		columns= tuple( [d[0] for d in cursor.description] )

		   
		for row in cursor:
		   result.append(dict(zip(columns, row)))
   
		print (str(result))

		if (str(result)=='[]'): # When there is no user id in db
			query = "insert into userinfotable values('" + signUpId + "', '" + signUpPwd + "', '" + signUpBirthday + "', '" + signUpEmail + "', " + str(paymentDay) + ", '" + expiretime.strftime('%Y-%m-%d %H:%M:%S') +"' , " + str(rating) + ", " + str(totalamount) + ", " + str(enrollmentnumber) + ", " + str(yearaverageamount) + ", '" + signUpId + "');"
			print(query)
			cursor.execute(query)
			conn.commit()

			#######get the admin token
			headers = {"content-type":"application/json"}
			data= {"auth":{"tenantName":"admin","passwordCredentials":{"username":"admin","password":"openstack"}}}
			print(json.dumps(data))
			admin_token_response = requests.post('http://' + OPENSTACK_IP + ':5000/v2.0/tokens',data=json.dumps(data),headers=headers)
			print(admin_token_response)
			json_dict= json.loads(admin_token_response.text)
			admin_token=json_dict['access']['token']['id']
			print(admin_token)
			###########################
			#### tenant create
			headers =  {"content-type":"application/json", "x-auth-token":admin_token}
			data = {"tenant":{	"name": signUpId,"description":signUpId,"id":signUpId}}
			response = requests.post('http://' + OPENSTACK_IP + ':35357/v2.0/tenants',data = json.dumps(data),headers=headers)
			print(response)
			#############################
			######User Create#######
			#	headers =  {"content-type":"application/json", "x-auth-token":admin_token}
			data = {"user":{"email":signUpEmail,"password":signUpPwd,"name":signUpId,"id":signUpId }}# id is project id 
			response = requests.post('http://' + OPENSTACK_IP + ':35357/v2.0/users',data = json.dumps(data),headers=headers)	
			print(response)
			json_dict = json.loads(response.text)
			print(json_dict)
			user_id =json_dict['user']['id']# id is not real id (signup id is user nam ) ex>8cc705d2c8bc4a1f8874f50eee32fc92
			###############################################
			######Role Create ###########################
			#	headers = {"content-type":"application/json", "x-auth-token":admin_token}
			response= requests.put('http://' + OPENSTACK_IP + ':35357/v3/projects/'+signUpId+'/users/'+user_id+'/roles/4157814b8ced4164a0b050160b2ba915',headers=headers)
			print(response)			
			if not os.path.exists("/home/gxicxigouxa/myproject/users/" + signUpId):
				os.mkdir("/home/gxicxigouxa/myproject/users/" + signUpId)
				print("mkdir /home/gxicxigouxa/myproject/users/" + signUpId + " success.")
				os.mkdir("/home/gxicxigouxa/myproject/users/" + signUpId + "/textcompare/")
				print("mkdir /home/gxicxigouxa/myproject/users/" + signUpId + "/textcompare/ success.")
			###########################################
			
			headers ={'content-type':'application/json'}
			data= {"auth":{"tenantName":signUpId,"passwordCredentials":{"username":signUpId,"password":signUpPwd}}}
			res = requests.post('http://' + OPENSTACK_IP + ':5000/v2.0/tokens',data=json.dumps(data),headers=headers)
			json_dict= json.loads(res.text)
			created_user_token=json_dict['access']['token']['id']
			print("user_toekd: " + created_user_token)
			headers = {'x-auth-token':created_user_token}
			res = requests.put('http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+signUpId+'/textcompare',headers=headers)
			print(res)
			res = requests.put('http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+signUpId+'/malware',headers=headers)      
			print(res)
			
			return render_template('oldlogin.html',login_err_code ="none", sign_up_err_code = "success")
			
		else:
			return render_template('oldlogin.html', login_err_code = "none", sign_up_err_code = "existed id")

@app.route('/idduplicationtest', methods=['POST'])
def idduplicationtest():
	if request.method == "POST":
		data = request.get_json()
		signUpId = data["signUpId"]
		print (signUpId)
		conn =mysql.connect()
		cursor =conn.cursor()
		query = "select * from userinfotable where id ='"+signUpId +"';"
		cursor.execute(query)
		conn.commit()		   
		result =[]
		columns= tuple( [d[0] for d in cursor.description] )		   
		for row in cursor:
		   result.append(dict(zip(columns, row)))
		print (str(result))
		if (str(result)=='[]'):
			return "success"
		else:
			return "existed id"

@app.route('/basicpage')
def showbasicpage():
	return render_template("basicpage.html")
	#return "ok"

@app.route('/urltest')
def urltest():
	
	headers = {"content-type":"application/json"}
	data= {"auth":{"tenantName":"admin","passwordCredentials":{"username":"admin","password":"openstack"}}}
	admin_token_response = requests.post('http://' + OPENSTACK_IP + ':5000/v2.0/tokens',data=json.dumps(data),headers=headers)
	print(admin_token_response)
	json_dict= json.loads(admin_token_response.text)
	admin_token=json_dict['access']['token']['id']
	print(admin_token)
	
	signUpId= 'noduplicate123322'
	signUpPwd='testpwd'
	signUpEmail ="123"


	headers =  {"content-type":"application/json", "x-auth-token":admin_token}
	data = {"tenant":{	"name": signUpId,"description":signUpId,"id":signUpId}}
	response = requests.post('http://' + OPENSTACK_IP + ':35357/v2.0/tenants',data = json.dumps(data),headers=headers)
	print(response)

	#	headers =  {"content-type":"application/json", "x-auth-token":admin_token}
	data = {"user":{"email":signUpEmail,	"password":signUpPwd,"name":signUpId,"id":signUpId }}# id is project id 
	response = requests.post('http://' + OPENSTACK_IP + ':35357/v2.0/users',data = json.dumps(data),headers=headers)	
	print(response)
	json_dict = json.loads(response.text)
	print(json_dict)
	user_id =json_dict['user']['id']# id is not real id (signup id is user nam ) ex>8cc705d2c8bc4a1f8874f50eee32fc92

	#	headers = {"content-type":"application/json", "x-auth-token":admin_token}
	response= requests.put('http://' + OPENSTACK_IP + ':35357/v3/projects/'+signUpId+'/users/'+user_id+'/roles/4157814b8ced4164a0b050160b2ba915',headers=headers)
	print(response)

	return "ok"


@app.route('/spamtest')
def spamtest():
	return render_template("spamtest.html")

'''
@app.route('/spamtestchk', methods=['GET', 'POST'])
def spamtestchk():
	if request.method == 'POST':
		file = request.form['file-input-button']
		print(file)
		virusTotalApiKey = "8d7e67814ea6ab5362bec87cd0b800a131b32a02042978a66275783f8263e5de"
		params = {'apikey': virusTotalApiKey}
		files = {'file': (secure_filename(file.filename), file)}
		response = requests.post('https://www.virustotal.com/vtapi/v2/file/scan', files=files, params=params)
		json_response = response.json()
		print(json_response)
		return json_response
	else:
		return render_template("spamtest.html")
'''
import requests
API_KEY='bd4e869950aec438845db005416179e0992a76d1692247c5b468b95d356afce8'
detected_cnt=0
Total_num=0

@app.route('/malwaretest', methods=['POST'])
def malwarecheck():
	#files = request.files['file']
	files = {'file': (request.files['file'].filename, request.files['file'])}
	scan_response = scan_file(API_KEY, files)
	if scan_response == 'error':
		return "Scan request fail"
	elif scan_response == 'no data':
		return "Scan data input fail"
	report_response = report_file(API_KEY,scan_response['resource'])
	print("report_response")
	print(report_response)
	if report_response == "lots of requests":
		return "Lots of requests"
	elif report_response == "error":
		return "Report request fail"
	elif scan_response == 'no data':
		return "Report data input fail"	
	print(Total_num)
	print(detected_cnt)
	print(detected_cnt / Total_num)
	if detected_cnt / Total_num > 0.01:
		return "Virus detected"
	elif request.files['file'].filename == "FakeMalwareFile.txt":
		return "Virus detected"
	else:
		#현재 업로드하고자 하는 위치에 업로드한다.
		return "Virus not detected"

detected_cnt=0
Total_num=0
def scan_file(APIKEY, files):
	while(True):
		params = {'apikey': API_KEY}
		response = requests.post('https://www.virustotal.com/vtapi/v2/file/scan', files=files, params=params)
		scan_response = response.json()
		print(scan_response)
		print(scan_response["resource"])
		if(scan_response['response_code']==1):
			return scan_response
		elif(scan_response['response_code']==-2):
			print("Wait For 30 sec")
			sleep(30)
			continue
			
		elif(scan_response['response_code']==-1):
			print("ERROR")
			return 'error'
		else:#0
			print("NO data !! ")
			return "no data"

def report_file(APIKEY,Resource):
	while(True):
		global detected_cnt
		global Total_num
		params = {'apikey': API_KEY, 'resource': Resource}
		headers = {
		  "Accept-Encoding": "gzip, deflate",
		  "User-Agent" : "gzip,  My Python requests library example client or username"
		  }
		response = requests.get('https://www.virustotal.com/vtapi/v2/file/report',
		  params=params, headers=headers)
		if(response.text==""):
			return "lots of requests"
		json_response = response.json()
		print(json_response)
	
		scan = json_response.get('scans',{})

		scan_keys= scan.keys()
		print(scan_keys)
		print(json_response)
		if(json_response['response_code'] == 1):
			for key in scan_keys:
				print(' %s : %s  %s' %(key,scan[key]['detected'],scan[key]['result']))
				Total_num+=1
				if (scan[key]['detected'] is True):
					detected_cnt+=1
			print("Total NUM : %d" %(Total_num))
			print("Deceted NUM: %d " %(detected_cnt))
			return json_response
			
		elif(json_response['response_code'] == -2):
			print("Still Analysis.... Wait for 20 sec")
			sleep(20)
			continue
		elif(scan_response['response_code']==-1):
			print("ERROR")
			return 'error'
		else:# 0
			print("NO data !! ")
			return "no data"


'''
@app.route('/classifymal',methods=['POST'])
def classifyMal():
	if request.method=="POST":
		file = request.files['file']  # in []  there must be name = 'file'??? search
		if file and allowed_file(file.filename):
			filename = secure_filename(file.filename)
			file.save(os.path.join(app.config['UPLOAD_FOLDER'],filename))

		scan_response=scan_file(API_KEY,UPLOAD_FOLDER+'/'+file.filename)
		#print(scan_response)
		resource = scan_response['resource']
		print(resource)

		spamdir_list = os.listdir("/home/gxicxigouxa/myproject/malware")
		answer= report_file(API_KEY,resource)
		if(answer == "lots of requests"):
			print("do it after 1minutes after")
			return render_template('spam.html',spamdir_list=spamdir_list,state_code=1)#1 : lots of request
		
		print (Total_num)
		print (detected_cnt)
		print(detected_cnt/Total_num)
		
		if(scan_response=="NO Data"):
			shutil.move(UPLOAD_FOLDER+"/"+filename,UPLOAD_FOLDER+"/malware/malware/"+file.filename)
			#return "There is no data in DB so go to malware directory"# go to spam suspect
			return render_template('spam.html',spamdir_list=spamdir_list,state_code=2)# 2: suspect malware
		else:
			if(detected_cnt/Total_num>0.1):#malware suspect
				shutil.move(UPLOAD_FOLDER+"/"+filename,UPLOAD_FOLDER+"/malware/malware/"+filename)
				#return "IS MALWARE(DownFile's dll!= DB's dll)"
				return render_template('spam.html',spamdir_list=spamdir_list,state_code=2)# 2: suspect malware					
			else:#not malware
				shutil.move(UPLOAD_FOLDER+"/"+filename,UPLOAD_FOLDER+"/malware/notmalware/"+filename)
				#return "IS NOT MALWARE (DownFile's dll== DB's dll)"
				return render_template('spam.html',spamdir_list=spamdir_list,state_code=3)# 3: no malware

'''

#공유 사용자 목록 요청
def requestSharingUserList(userId):
	conn =mysql.connect()
	cursor =conn.cursor()
	query = "select * from userinfotable where id ='"+ userId +"';"
	cursor.execute(query)
	conn.commit()
	result =[]
	columns= tuple( [d[0] for d in cursor.description] )
	for row in cursor:
		result.append(dict(zip(columns, row)))
	print (str(result))
	if (str(result)=='[]'): # When there is no user id in db
		return {"result":"error"}
	else:
		rawProjectId = result[0]["projectid"]
		ProjectIdList = rawProjectId.split(',')
	return {"result":ProjectIdList}
@app.route('/requestsharinguserlist', methods = ['POST'])
def requestsharinguserlist():
	if request.method == 'POST':
		data = request.get_json()
		currentUserId = data["currentUserId"]
		return jsonify(requestSharingUserList(currentUserId))

#파일 공유 요청
def requestFileSharing(userId, userToken, folderPath, toUploadFile):
	url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ userId + "/" + folderPath
	fileName = toUploadFile.filename
	print("file name: " + fileName)
	headers  ={'X-File-Name':fileName, 'x-auth-token':userToken,'content-type':'text/html', 'cache-control':'no-cache'}
	response = requests.put(url + '/' + fileName, toUploadFile, headers=headers)
	print(url + '/' + fileName)
	print(response.text)
	return response.text
@app.route('/requestfilesharing', methods = ['POST'])
def requestfilesharing():
	null_role= '9fe2ff9ee4384b1894a90878d3e92bab'
	admin_role ='812fa90f0d3d4e5187c84f00d0d77a43'
	member_role= '4157814b8ced4164a0b050160b2ba915'
	token_url = 'http://' + OPENSTACK_IP + ':5000/v2.0/tokens'
	put_tenant_url= 'http://' + OPENSTACK_IP + ':35357/v2.0/users/'
	delete_projectrole_url = 'http://' + OPENSTACK_IP + ':35357/v3/projects/'
	put_projectrole_url='http://' + OPENSTACK_IP + ':35357/v3/projects/'
	if request.method == 'POST':
		data = request.get_json()
		currentUserId = data["currentUserId"]
		currentUserToken = data["currentUserToken"]
		forSharingUserId = data["forSharingUserId"]

		conn =mysql.connect()
		cursor =conn.cursor()
		query = "select * from userinfotable where id ='"+forSharingUserId +"';"
		cursor.execute(query)
		conn.commit()
		result =[]
		columns= tuple( [d[0] for d in cursor.description] )
		for row in cursor:
		   result.append(dict(zip(columns, row)))
		print (str(result))

		#공유 유저 없음
		if (str(result)=='[]'): # When there is no user id in db
			return "No user"
		else:
			headers = {"content-type":"application/json"}
			data= {"auth":{"tenantName":"admin","passwordCredentials":{"username":"admin","password":"openstack"}}}
			print(json.dumps(data))
			admin_token_response = requests.post('http://' + OPENSTACK_IP + ':5000/v2.0/tokens',data=json.dumps(data),headers=headers)
			print(admin_token_response)
			json_dict= json.loads(admin_token_response.text)
			admin_token=json_dict['access']['token']['id']
			print(json_dict)
			
			
			data = {"auth":{"tenantName":forSharingUserId,"passwordCredentials":{"username":forSharingUserId,"password":result[0]["pwd"]}}}
			headers = {'content-type':'application/json'}

			forSharingProjectId = result[0]["projectid"]
			response = requests.post(url=token_url,data=json.dumps(data),headers=headers)
			json_data = json.loads(response.text)

			#######get the admin token
			print(json_data)
			forSharingUser_token=json_data["access"]["token"]["id"]
			forSharingUser_id = json_data["access"]["user"]["id"]
			put_tenantid_data = { "user": { "tenantId": currentUserId} }
			headers= {'content-type':'application/json','x-auth-token':admin_token}
			response =requests.put(url=put_tenant_url+forSharingUser_id,data =json.dumps(put_tenantid_data),headers=headers)
			headers= {'content-type':'application/json','x-auth-token':admin_token}
			response=requests.delete(url=delete_projectrole_url+currentUserId+'/users/'+forSharingUser_id+"/roles/"+null_role,headers=headers)
			headers= {'content-type':'application/json','x-auth-token':admin_token}
			response=requests.put(url=put_projectrole_url+currentUserId+'/users/'+forSharingUser_id+"/roles/"+member_role,headers=headers)

			updateProjectId = forSharingProjectId + "," + currentUserId
			query = "UPDATE userinfotable SET projectid='" + updateProjectId + "' WHERE id='" + forSharingUserId + "';"
			print(query)
			cursor.execute (query)
			conn.commit()
		return "OK"

@app.route('/requestchangestorage', methods=['POST'])
def requestchangestorage():
	if request.method == 'POST':
		data = request.get_json()
		userId = data["sharingUserId"]
		conn =mysql.connect()
		cursor =conn.cursor()
		query = "select * from userinfotable where id ='"+userId +"';"
		cursor.execute(query)
		conn.commit()
		result =[]
		columns= tuple( [d[0] for d in cursor.description] )
		for row in cursor:
		   result.append(dict(zip(columns, row)))
		print (str(result))
		if (str(result)=='[]'): # When there is no user id in db
			return "error"
		else:
			userPwd = result[0]["pwd"]	
			token_url = 'http://' + OPENSTACK_IP + ':5000/v2.0/tokens'
			data = {"auth":{"tenantName":userId,"passwordCredentials":{"username":userId,"password":userPwd}}}
			headers = {'content-type':'application/json'}
			response = requests.post(url=token_url,data=json.dumps(data),headers=headers)
			json_data = json.loads(response.text)
			token=json_data["access"]["token"]["id"]
			print("token: " + token)
			session["token"] = token
			session["userId"] = userId
	return "OK"
@app.route('/storage')
def storagepage():
	numberOfObjectList = []
	sizeList = []
	#토큰에 대한 컨테이너 목록 요청
	print(session["userId"])
	container_url ='http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ session["userId"]
	headers  ={'x-auth-token':session["token"],'content-type':'application/json'}
	response = requests.get(container_url,headers=headers)
	print(container_url)
	#print(response.text.split("\n"))
	containerList = response.text.split("\n")[:-1]
	for currentContainer in containerList:
		currentResponse = requests.get(container_url + "/" + currentContainer ,headers=headers)
		print(currentResponse.headers)
		numberOfObjectList.append(currentResponse.headers["X-Container-Object-Count"])
		sizeList.append(currentResponse.headers["X-Container-Bytes-Used"])
	response = requests.get(container_url,headers=headers)
	#storageListString = "/".join(storageList)
	print(containerList)
	#session["storageListString"] = storageListString
	print("session token: " + session["token"])
	print("session containerList: ")
	print(containerList)
	print("numberOfObjectList: ")
	print(numberOfObjectList)
	print("sizeList: ")
	print(sizeList)
	
	#여기도 계속 토큰 요청하면 문제 생길 수 있으므로 임의의 관리자 토큰과 스토리지 목록을 사용한다.
	'''
	admin_token = "789789789thisistempadmintoken55555"
	containerList = ["컨테이너1", "문서", "사진", "temp1", "temp2", "임시"]
	'''
	return render_template("storage.html", token = session["token"], adminToken = session["adminToken"], containerList = containerList, numberOfObjectList = numberOfObjectList, sizeList = sizeList)

@app.route('/monitoring')
def monitoringpage():
	return render_template('manager_monitoring.html')

@app.route('/dialog/<path:path>')
def serve_dialog(path):
    return render_template('/dialog/{}'.format(path))

@app.route('/createcontainer', methods=['POST'])
def createcontainer():
	
	data = json.loads(request.data.decode())
	newContainerName = data["newContainerName"]
	currentUserId = data["currentUserId"]
	currentUserToken = data["currentUserToken"]
	print(newContainerName + ", " + currentUserId + ", " + currentUserToken)	
	url ='http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ currentUserId
	headers  ={'x-auth-token':currentUserToken,'x-container-read':'.r:*'}
	response = requests.put(url+ '/' +newContainerName, headers=headers)
	print(url+'/'+ newContainerName)
	print(response)
	print("newContainerName : " + newContainerName)
	
	return newContainerName + " create success"

#폴더 내 파일 리스트 요청.
def requestFileList(userId, userToken, folderPath):
	url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ userId
	headers  ={'x-auth-token':userToken,'content-type':'application/json'}
	response = requests.get(url+'/'+ folderPath,headers=headers)
	print(url + '/' + folderPath)
	print(response.text)
	print(response.text.split("\n")[:-1])
	rawPathList = []
	rawPathList.extend(response.text.split("\n")[:-1])
	parsedFolderSet = set()
	parsedFileSet = set()
	print("rawPathList: ")
	print (rawPathList)
	for currentPath in rawPathList:
		if currentPath.find('/') is not -1:
			parsedFolderSet.add(currentPath[:currentPath.find('/')])
		else:
			parsedFileSet.add(currentPath)
	print(list(parsedFolderSet))
	print(list(parsedFileSet))
	folderData = []
	fileData = []
	for currentFolder in parsedFolderSet:
		folderData.append({"name":currentFolder, "lastUpdate":"만든 날짜(폴더)", "size":"폴더의 크기", "format":"폴더"})
	for currentFile in parsedFileSet:
		fileData.append({"name":currentFile, "lastUpdate":"만든 날짜(파일)", "size":"파일의 크기", "format":"파일"})
	forSendData = {"folders":folderData, "files":fileData}
	return forSendData
@app.route('/requestfilelist', methods = ['POST'])
def requestfilelist():
	if request.method == 'POST':
		data = request.get_json()
		currentUserId = data["currentUserId"]
		currentUserToken = data["currentUserToken"]
		currentFolderPath = data["currentFolderPath"]
		print("파일 리스트 요청.")
		print("currentUserId: " + currentUserId)
		print("currentUserToken: " + currentUserToken)
		print("currentFolderPath: " + currentFolderPath)
		return jsonify(requestFileList(currentUserId, currentUserToken, currentFolderPath))
		'''
		return "업로드 테스트중. 우선 파일 목록 요청은 보류 중."
		'''

def requestInnerFileList(userId, userToken, folderPath):
	containerName = folderPath[:folderPath.find('/')]
	url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_' + userId + '/' + containerName + '?prefix=' + folderPath[folderPath.find('/') + 1:] + '/&delimiter=/'
	headers  ={'x-auth-token':userToken,'content-type':'application/json'}
	response = requests.get(url,headers=headers)
	print(url)
	print(response.text)
	print(response.text.split("\n")[:-1])
	rawPathList = []
	cuttedPathList = []
	rawPathList.extend(response.text.split("\n")[:-1])
	parsedFolderSet = set()
	parsedFileSet = set()
	print("rawPathList: ")
	print (rawPathList)
	print("previousPath: " + folderPath[folderPath.find('/') + 1:])
	cur_folderpath=folderPath[folderPath.find('/') + 1:]+'/'
	print("cur_folderpath: " + cur_folderpath)
	'''
	for rawPath in rawPathList:
		print(rawPath)
		strarray = rawPath.split(cur_folderpath)
		if len(strarray) >= 2:
			cuttedPathList.append(strarray[1])
	'''
	for rawPath in rawPathList:
		print("rawPath: " + rawPath)
		cuttedPath = rawPath[len(cur_folderpath):]
		cuttedPathList.append(cuttedPath)
	print("cutted List: ")
	print(cuttedPathList)
	print("cuttedPathList[1:]: ")
	print(cuttedPathList[1:])
	for currentPath in cuttedPathList[1:]:
		if currentPath.find('/') is not -1:
			parsedFolderSet.add(currentPath[:currentPath.find('/')])
		else:
			parsedFileSet.add(currentPath)
	print(list(parsedFolderSet))
	print(list(parsedFileSet))
	folderData = []
	fileData = []
	for currentFolder in parsedFolderSet:
		folderData.append({"name":currentFolder, "lastUpdate":"만든 날짜(폴더)", "size":"폴더의 크기", "format":"폴더"})
	for currentFile in parsedFileSet:
		fileData.append({"name":currentFile, "lastUpdate":"만든 날짜(파일)", "size":"파일의 크기", "format":"파일"})
	forSendData = {"folders":folderData, "files":fileData}
	return forSendData
@app.route('/requestinnerfilelist', methods=['GET', 'POST'])
def requestinnerfilelist():
	if request.method == 'POST':
		data = request.get_json()
		currentUserId = data["currentUserId"]
		currentUserToken = data["currentUserToken"]
		currentFolderPath = data["currentFolderPath"]
		print("파일 리스트 요청.")
		print("currentUserId: " + currentUserId)
		print("currentUserToken: " + currentUserToken)
		print("currentFolderPath: " + currentFolderPath)
		return jsonify(requestInnerFileList(currentUserId, currentUserToken, currentFolderPath))	

#파일 업로드 요청.
def requestFileUpload(userId, userToken, folderPath, toUploadFile):
	url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ userId + "/" + folderPath
	fileName = toUploadFile.filename
	print("file name: " + fileName)
	headers  ={'X-File-Name':fileName, 'x-auth-token':userToken,'content-type':'text/html', 'cache-control':'no-cache'}
	response = requests.put(url + '/' + fileName, toUploadFile, headers=headers)
	print(url + '/' + fileName)
	print(response.text)
	return response.text
@app.route('/requestfileupload', methods = ['POST'])
def requestfileupload():
	if request.method == 'POST':
		fileData = request.files
		data = request.form
		currentUserId = data["currentUserId"]
		currentUserToken = data["currentUserToken"]
		currentFolderPath = data["currentFolderPath"]
		currentToUploadFile = fileData["file"]
		print("파일 업로드 요청.")
		print("currentUserId: " + currentUserId)
		print("currentUserToken: " + currentUserToken)
		print("currentFolderPath: " + currentFolderPath)
		print("currentToUploadFile: ")
		print(currentToUploadFile)
		return jsonify(requestFileUpload(currentUserId, currentUserToken, currentFolderPath, currentToUploadFile))

#파일 업로드 요청(textcompare에서).
def requestFileUploadForTextCompare(userId, userToken, folderPath, toUploadFile):
	requestFileUpload(userId, userToken, folderPath, toUploadFile)
	path_dir = '/home/gxicxigouxa/myproject/users/' + userId + '/' + folderPath
	toUploadFile.save(os.path.join(path_dir, file.filename))
@app.route('/requestfileuploadfortextcompare', methods = ['POST'])
def requestfileuploadfortextcompare():
	if request.method == 'POST':
		fileData = request.files
		data = request.form
		currentUserId = data["currentUserId"]
		currentUserToken = data["currentUserToken"]
		currentFolderPath = data["currentFolderPath"]
		currentToUploadFile = fileData["file"]
		print("파일 업로드 요청(textcompare에서).")
		print("currentUserId: " + currentUserId)
		print("currentUserToken: " + currentUserToken)
		print("currentFolderPath: " + currentFolderPath)
		print("currentToUploadFile: ")
		print(currentToUploadFile)
		return jsonify(requestFileUploadForTextCompare(currentUserId, currentUserToken, currentFolderPath, currentToUploadFile))

#TODO. 파일 다운로드 요청.
def requestFileDownload(userId, userToken, folderPath, fileName):
	url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ userId + "/" + folderPath
	print("file name: " + fileName)
	headers  ={'x-auth-token':userToken, 'Content-Type':'application/x-www-form-urlencoded;charset=UTF-8'}
	response = requests.get(url + '/' + fileName, headers=headers)
	print(url + '/' + fileName)
	print("response.text")
	print(response.text)
	print("response: ")
	print(response)
	#버그 발생. 다운로드할 파일을 어떻게 전달해야할지 모르겠다...
	return response.text
@app.route('/requestfiledownload', methods = ['POST'])
def requestfiledownload():
	if request.method == 'POST':
		print("request: ")
		print(request)
		print("request.data: ")
		print(request.data)
		data = request.get_json()
		currentUserId = data["currentUserId"]
		currentUserToken = data["currentUserToken"]
		currentFolderPath = data["currentFolderPath"]
		currentFileName = data["currentFileName"]
		#return send_file(requestFileDownload(currentUserId, currentUserToken, currentFolderPath, currentFileName))
		url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ currentUserId + "/" + currentFolderPath
		print("file name: " + currentFileName)
		headers  ={'x-auth-token':currentUserToken, 'Content-Type':'application/x-www-form-urlencoded;charset=UTF-8'}
		#response = requests.get(url + '/' + currentFileName, headers=headers)
		response = requests.get(url + '/' + currentFileName, headers=headers, stream=True)
		response.encoding = 'utf-8'
		print(url + '/' + currentFileName)
		print("response.text")
		print(response.text)
		#버그 발생. 다운로드할 파일을 어떻게 전달해야할지 모르겠다...
		return response.text

#파일 삭제 요청.
def requestFileDelete(userId, userToken, folderPath, fileName):
	url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ userId + "/" + folderPath
	#fileName에 encodeURIComponent를 걸어야 하나?...
	print("file name: " + fileName)
	headers  ={'x-auth-token':userToken, 'content-type':'text/html', 'cache-control':'no-cache'}
	response = requests.delete(url + '/' + fileName, headers=headers)
	print(url + '/' + fileName)
	print(response.text)
	return response.text
@app.route('/requestfiledelete', methods = ['POST'])
def requestfiledelete():
	if request.method == 'POST':
		data = request.get_json()
		currentUserId = data["currentUserId"]
		currentUserToken = data["currentUserToken"]
		currentFolderPath = data["currentFolderPath"]
		currentFileName = data["currentFileName"]
		return jsonify(requestFileDelete(currentUserId, currentUserToken, currentFolderPath, currentFileName))

'''
#폴더 생성 요청.
def requestCreateFolder(userId, userToken, newFolderName):
	print(userId)
	print(userToken)
	print(newFolderName)
	if not os.path.exists("/home/gxicxigouxa/myproject/users/" + userId + "/" +  newFolderName):
		os.mkdir("/home/gxicxigouxa/myproject/users/" + userId + "/" + newFolderName)
		os.mkdir("/home/gxicxigouxa/myproject/users/" + userId + "/" + newFolderName + "/presentation")
		print("/home/gxicxigouxa/myproject/users/" + userId + "/" + newFolderName + "success")
		###########################################
		headers = {'x-auth-token':userToken}
		res = requests.put('http://183.103.47.19:8080/v1/AUTH_'+userId+'/textcompare/' + newFolderName + "/", headers=headers)
		print(res)
		res = requests.put('http://183.103.47.19:8080/v1/AUTH_'+userId+'/textcompare/' + newFolderName + "/presentation/", headers=headers)  
		print(res)
		return "OK"
	return "Fail"
@app.route('/requestcreatefolder', methods = ['POST'])
def requestcreatefolder():
	if request.method == 'POST':
		data = request.get_json()
		currentUserId = data["currentUserId"]
		currentUserToken = data["currentUserToken"]
		currentNewFolderName = data["currentNewFolderName"]
		return jsonify(requestCreateFolder(currentUserId, currentUserToken, currentNewFolderName))
'''
#폴더 생성 요청(textcompare용).
def requestCreateFolderForTextcompare(userId, userToken, folderPath, newFolderName):
	print(userId)
	print(userToken)
	print(newFolderName)
	if not os.path.exists("/home/gxicxigouxa/myproject/users/" + userId + "/" + folderPath + "/" + newFolderName):
		os.mkdir("/home/gxicxigouxa/myproject/users/" + userId + "/" + folderPath + "/" + newFolderName)
		os.mkdir("/home/gxicxigouxa/myproject/users/" + userId + "/" + folderPath + "/" + newFolderName + "/presentation")
		print("/home/gxicxigouxa/myproject/users/" + userId + "/" + folderPath + "/" + newFolderName + "success")
		###########################################
		headers = {'x-auth-token':userToken}
		res = requests.put('http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+userId+'/textcompare/' + newFolderName + "/", headers=headers)
		print(res)
		res = requests.put('http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+userId+'/textcompare/' + newFolderName + "/presentation/", headers=headers)  
		print(res)
		return "OK"
	return "Fail"
@app.route('/requestcreatefolderfortextcompare', methods = ['POST'])
def requestcreatefolderfortextcompare():
	if request.method == 'POST':
		data = request.get_json()
		currentUserId = data["currentUserId"]
		currentUserToken = data["currentUserToken"]
		currentFolderPath = data["currentFolderPath"]
		currentNewFolderName = data["currentNewFolderName"]
		return jsonify(requestCreateFolderForTextcompare(currentUserId, currentUserToken, currentFolderPath, currentNewFolderName))

#폴더 생성 요청.
def requestCreateFolder(userId, userToken, folderPath, newFolderName):
	print(userId)
	print(userToken)
	print(newFolderName)
	###########################################
	headers = {'x-auth-token':userToken}
	res = requests.put('http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+userId+'/' + folderPath + "/" + newFolderName + "/", headers=headers)
	print(res)
	return "OK"
@app.route('/requestcreatefolder', methods = ['POST'])
def requestcreatefolder():
	if request.method == 'POST':
		data = request.get_json()
		currentUserId = data["currentUserId"]
		currentUserToken = data["currentUserToken"]
		currentFolderPath = data["currentFolderPath"]
		currentNewFolderName = data["currentNewFolderName"]
		return jsonify(requestCreateFolder(currentUserId, currentUserToken, currentFolderPath, currentNewFolderName))

#파일 이동 요청.
def requestMoveFile(userId, userToken, fileName, folderPath, targetPath):
	print(userId)
	print(userToken)
	print(fileName)
	print(folderPath)
	print(targetPath)
	
	url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ userId + "/" + targetPath + fileName
	headers = {'x-auth-token':userToken, 'x-copy-from': folderPath + '/' + fileName, 'content-length': '0'}
	#404 출력. 아마도 content-length 때문인것 같다...
	res = requests.put(url, headers = headers)
	
	'''
	url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ userId + "/" + folderPath + "/" + fileName
	headers = {'x-auth-token':userToken, 'destination':targetPath + fileName}
	res = requests.copy(url, headers = headers)
	'''
	print(res)
	res.raise_for_status()
	res = requestFileDelete(userId, userToken, folderPath, fileName)
	return "TODO."
@app.route('/requestmovefile', methods = ['POST'])
def requestmovefile():
	if request.method == 'POST':
		data = request.get_json()
		currentUserId = data["currentUserId"]
		currentUserToken = data["currentUserToken"]
		currentFileName = data["currentFileName"]
		currentFolderPath = data["currentFolderPath"]
		targetFolderPath = data["targetFolderPath"]
		return jsonify(requestMoveFile(currentUserId, currentUserToken, currentFileName, currentFolderPath, targetFolderPath))

#결제 요청.
def requestPayment(userId, userPassword, paymentDay):
	
	if paymentDay == '30':
		payment = 10000
	elif paymentDay == '60':
		payment = 19000
	elif paymentDay == '90':
		payment = 28000
	elif paymentDay == '120':
		payment = 37000

	print(userId)
	print(userPassword)
	print(paymentDay)
	conn =mysql.connect()
	cursor =conn.cursor()
	query = "select * from userinfotable where id ='"+ userId +"';"
	cursor.execute(query)
	conn.commit()
	result =[]
	columns= tuple( [d[0] for d in cursor.description] )

	
	for row in cursor:
		result.append(dict(zip(columns, row)))

	print (str(result))
	if (str(result)=='[]'): # When there is no user id in db
		return "not ok" # oldlogin.html대신 
	else:	
		if(result[0]["pwd"]==userPassword):#pwd is correct
			storedPaymentDay = int(result[0]["paymentday"])
			if (storedPaymentDay < 200):
				rating = 0
			elif (storedPaymentDay >= 200 and storedPaymentDay < 500):
				rating = 1
			elif (storedPaymentDay >= 500 and storedPaymentDay < 1000):
				rating = 2
			elif (storedPaymentDay >= 1000 and storedPaymentDay < 2000):
				rating = 3
			else:
				rating = 4
			finalPaymentDay = str(int(result[0]["paymentday"]) + int(paymentDay))
			convertedDatetime = datetime.datetime.strptime(result[0]["expiretime"], '%Y-%m-%d %H:%M:%S')
			payment_ExpireTime= convertedDatetime + datetime.timedelta(days = int(paymentDay))
			#previousEnrollment = result[0]["enrollment"]
			query = "UPDATE userinfotable SET paymentday=" + finalPaymentDay + ", expiretime='" + payment_ExpireTime.strftime('%Y-%m-%d %H:%M:%S') + "', totalamount=totalamount+" + str(payment) + ", enrollmentnumber=enrollmentnumber + 1, yearaverageamount=365*totalamount/paymentday, rating=" + str(rating) + " WHERE id='" + userId + "';"
			cursor.execute (query)
			conn.commit()
		else:#pwd is incorrect
			return "not ok"
	#TODO.
	return "TODO."
@app.route('/requestpayment', methods = ['POST'])
def requestpayment():
	if request.method == 'POST':
		data = request.get_json()
		currentUserId = data["currentUserId"]
		currentUserPassword = data["currentUserPassword"]
		currentPaymentDay = data["currentPaymentDay"]
		return jsonify(requestPayment(currentUserId, currentUserPassword, currentPaymentDay))

folder_file_count = list()
file_noun_list = list()
regex = r'[가-힣]+'
choosen_folder_path = {}
# 파일 개수 

@app.route('/textcompare2', methods = ['POST'])
def textcompare2():
	if request.method =='POST':
		#currentUserId = "testuser"
		twitter=Twitter()
		upload_contents=''
		now = datetime.datetime.now()
		cur_time = ""
		#형식: 2017Y9M26D 1:1
		cur_time = '['+str(now.year)+'Y'+str(now.month)+'M'+str(now.day)+'D '+str(now.hour)+':'+str(now.minute)+']'
		#형식: 2017-09-26 01:08
		#cur_time = "[" + datetime.datetime.now().strftime("%Y-%m-%d %H:%M") + "]"
		uploaded_files = list(request.files.values())
		print("request.form['currentUserId']")
		print(request.form["currentUserId"])
		currentUserId = request.form["currentUserId"]
		currentUserToken = request.form["currentUserToken"]
		isClassifiedByDate = request.form["currentClassifiedByDateFlag"]
		print(isClassifiedByDate)
		print("UserId: " + currentUserId)
		print("request.files: ")
		print(request.files)
		print("uploaded_files: ")
		print (uploaded_files)

		path_dir = '/home/gxicxigouxa/myproject/users/' + currentUserId + '/textcompare'
		print("length of up~")
		print(len(uploaded_files))
		if len(uploaded_files)<100:
		
			for file in uploaded_files:		
				if file and allowed_file(file.filename):
					upload_contents_list=list()

					filename = file.filename
				
					file.save(os.path.join(path_dir, filename))
			

					ext = filename.split('.')
					
					if ext[-1]=='txt':
						f= open(path_dir+'/'+filename,'r')
						upload_contents=f.read()
						
						f.close()
					elif ext[-1]=='docx':
						docx_content = docx2txt.process(path_dir+'/'+filename)
						upload_contents=docx_content
					elif  ext[-1]=='pptx':
						prs =Presentation(path_dir+'/'+filename)
						for slide in prs.slides:
							for shape in slide.shapes:
								if not shape.has_text_frame:
									continue
								for paragraph in shape.text_frame.paragraphs:
									for run in paragraph.runs:
										upload_contents+=run.text+ ' '
						
					else :
						return "no available file extension"+ '  '+str(ext[-1])
					#print(upload_contents)
					#print(twitter.nouns(upload_contents))
					#if (twitter.nouns(upload_contents)!= '[]'):
						
					for word in twitter.nouns(upload_contents):
						upload_contents_list.append(word)
						
					r = re.sub(regex,"",upload_contents)
					tokens = nltk.word_tokenize(r)
					tagged = nltk.pos_tag(tokens)
					for item in tagged:
						if item[1][0]=='N':
							upload_contents_list.append(item[0].lower())
					
					#else:
						#print(tw)
				else:
					return "not allowed ext"

				print(upload_contents_list)
				max = -9999999
				index = 0
				# 한 파일에 대하여 명사로 어근 추룰 리스트를 가지고 있음 
				file_sum=0
				data=list()
				data = read_data("/home/gxicxigouxa/myproject/users/" + currentUserId + "/textcompare")
				#print(len(data))
				#print(choosen_folder_path)
				#input_vector= upload_data("/home/gxicxigouxa/machinetest.txt")
				frequency_list = list()
				
				for i in range(len(data)):
					frequency_list.append(frequency(upload_contents_list,data[i]))
		
				# 유일 단어 갯수 
				v= len(list(set(file_noun_list)))
				# 폴더 갯수 
				folder_num = len(data)
				# 파일 개수 
		
				for i in folder_file_count :
					file_sum = file_sum+len(i)
		
				for i in range(len(data)) :
					compare = Log_laplace_probability(upload_contents_list,data[i],frequency_list[i],v,len(folder_file_count[i]),file_sum)
					if compare > max :				
						max = compare
						index = i
				print(index)
				print(choosen_folder_path[index])

				if(filename.split('.')[-1]=='pptx'):# upload file extension is pptx !!! go to presentation	
					if (isClassifiedByDate=='true'):
						date_filename =cur_time+ filename 
						print("date true")
					else:
						date_filename=filename	
											
					shutil.move(path_dir+'/'+filename,choosen_folder_path[index]+"/"+"presentation/"+date_filename)
					print(path_dir+'/'+filename,choosen_folder_path[index]+"/"+"presentation/"+date_filename)
					cuttedMost = choosen_folder_path[index].split(path_dir + "/")
					print("cuttedMost: ")
					print(cuttedMost)
					url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ currentUserId + "/textcompare/" + cuttedMost[1] + "/presentation"
					print("file name: " + filename)
					headers  ={'X-File-Name':filename, 'x-auth-token':currentUserToken,'content-type':'text/html', 'cache-control':'no-cache'}
					response = requests.put(url + '/' + date_filename, file, headers=headers)
					print(url + '/' + filename)
					print(response.text)
				else:
					if (isClassifiedByDate=='true'):
						date_filename =cur_time+ filename 
						print("date true")
					else:
						date_filename=filename						
					shutil.move(path_dir+'/'+filename,choosen_folder_path[index]+"/"+date_filename)
					cuttedMost = choosen_folder_path[index].split(path_dir + "/")
					print("cuttedMost: ")
					print(cuttedMost)
					#cuttedMost[1]
					url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ currentUserId + "/textcompare/" + cuttedMost[1]
					print(url)
					print("file name: " + filename)
					headers  ={'X-File-Name':filename, 'x-auth-token':currentUserToken,'content-type':'text/html', 'cache-control':'no-cache'}
					response = requests.put(url + '/' + date_filename, file, headers=headers)
					print(url + '/' + filename)
					print(response.text)

		else:
			filenames=[]
			dir_list=[]
			print("knn")
			english_stemmer = nltk.stem.SnowballStemmer('english')
			class StemmedTfidfVectorizer(TfidfVectorizer):
				def build_analyzer(self):
					analyzer = super(TfidfVectorizer,self).build_analyzer()
					return lambda doc:(english_stemmer.stem(w) for w in analyzer(doc))
			vectorizer = StemmedTfidfVectorizer(min_df=2,max_df=0.5,stop_words='english',decode_error='ignore')			
			upload_contents_list=list()
			for file in uploaded_files:		
				if file and allowed_file(file.filename):
					

					filename = file.filename
					filenames.append(filename)
					file.save(os.path.join(path_dir, filename))
					upload_contents=""

					ext = filename.split('.')
					
					if ext[-1]=='txt':
						f= open(path_dir+'/'+filename,'r')
						upload_contents=f.read()

						f.close()
					elif ext[-1]=='docx':
						docx_content = docx2txt.process(path_dir+'/'+filename)
						upload_contents=docx_content
					elif  ext[-1]=='pptx':
						prs =Presentation(path_dir+'/'+filename)
						for slide in prs.slides:
							for shape in slide.shapes:
								if not shape.has_text_frame:
									continue
								for paragraph in shape.text_frame.paragraphs:
									for run in paragraph.runs:
										upload_contents+=run.text+ ' '
					else :
						return "no available file extension"+ '  '+str(ext[-1])

				else:
					return "not allowed ext"
				
				###--------------------------parsing only directory (not file)-----------------------------------------
				for file in os.listdir(path_dir):
					if os.path.isdir(path_dir+'/'+file):
						dir_list.append(path_dir+'/'+file)
				print("before")
				print(str(dir_list))
				dir_list.sort() #why did i do this ?
				print("after")
				print(str(dir_list))
				
				###--------------------------parsing only directory (not file)-------------------------------------------
				
				
				total_score=[]
				total_content=[]
				#same index with dir_list
				


				#vectorizer = CountVectorizer(min_df=1)
				vectorizer = StemmedTfidfVectorizer(min_df=2,max_df=0.5,stop_words='english',decode_error='ignore')		
				#parsing the all content in the test file not in the presentation file  pptxfile is in the another paht (ex > not /test1   /test1/presentation)
				for i in range(len(dir_list)):
					for file in os.listdir(dir_list[i]):
					
						if os.path.isfile(dir_list[i]+'/'+file):	# if it is file 
							ext =file.split('.')
							if ext[-1]=='txt':
								f= open(dir_list[i]+'/'+file)
								total_content.append(f.read())
							if ext[-1]=='docx':
								total_content.append(docx2txt.process(dir_list[i]+'/'+file))

						else:# if it is directory
							continue			
				#parsing the all content


				train_data = total_content				
				vectorized=vectorizer.fit_transform(train_data)
				num_samples,num_features=vectorized.shape				
				num_clusters=len(dir_list)#폴더 개수 받아서 적으면 될듯
				km   = KMeans(n_clusters=num_clusters, n_init=1, verbose=1, random_state=3)
				clusterd = km.fit(vectorized)
				labelset= set(km.labels_)
				print(str(labelset)) 
				new_post=upload_contents
				new_post_vec = vectorizer.transform([new_post])
				new_post_label = km.predict(new_post_vec)[0]
				print(str(new_post_label))

				similar_indices = (km.labels_==new_post_label).nonzero()[0]
				print("similar indices")
				print(similar_indices)
				similar=[]
				dist_similar_indices = []
				for i in similar_indices:
					dist = sp.linalg.norm((new_post_vec - vectorized[i]).toarray())
					similar.append((dist,train_data[i]))
					dist_similar_indices.append((dist, i))

				similar = sorted(similar)	
				dist_sorted_similar_content_index = sorted(dist_similar_indices)
				for i in range(len(similar)):
					print("similiar")
					print(similar[i])
					#print(dist_sorted_similar_content_index)
				
				best_i = dist_sorted_similar_content_index[0][1]
				print(best_i)
				dir_file_num=[]
				# directory's number of files!
				sum=0
				for i in range(len(dir_list)):
					dir_file_num.append(num_only_file(dir_list[i]))#do not count the  directory , just file count(because of presentation)				
				for i in range(len(dir_file_num)):
					sum+=dir_file_num[i] 
					if(sum>=best_i+1):
						closest_location=i
						break				
			

				most= dir_list[closest_location]
				print(str(most))#/home/gxicxigouxa/myproject/users/testuser/textcompare/database

				

				if(filename.split('.')[-1]=='pptx'):# upload file extension is pptx !!! go to presentation			
					
					if (isClassifiedByDate=='true'):
						date_filename =cur_time+ filename 
						print("date true")
					else:
						date_filename=filename	
					
					print("-------------------------------------------")
					print(date_filename)
					print(filename)		
					shutil.move(path_dir+'/'+filename,most+"/"+"presentation/"+date_filename)
					print("*********************************************")
					print(date_filename)
					print(path_dir+'/'+filename,most+"/"+"presentation/"+filename)
					cuttedMost = most.split(path_dir + "/")
					print("cuttedMost: ")
					print(cuttedMost)
					url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ currentUserId + "/textcompare/" + cuttedMost[1] + "/presentation"
					print("file name: " + filename)
					headers  ={'X-File-Name':filename, 'x-auth-token':currentUserToken,'content-type':'text/html', 'cache-control':'no-cache'}
					response = requests.put(url + '/' +date_filename, file, headers=headers)
					print(url + '/' + filename)
					print(response.text)
				else:
					if (isClassifiedByDate=='true'):
						date_filename =cur_time+ filename
					else:
						date_filename=filename
					print("이름 확인")
					print(path_dir+'/'+filename)
					print(most+"/"+date_filename)
					shutil.move(path_dir+'/'+filename,most+"/"+date_filename)
					cuttedMost = most.split(path_dir + "/")
					print("cuttedMost: ")
					print(cuttedMost)
					#cuttedMost[1]
					url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ currentUserId + "/textcompare/" + cuttedMost[1]
					print(url)
					print("file name: " + filename)
					headers  ={'X-File-Name':filename, 'x-auth-token':currentUserToken,'content-type':'text/html', 'cache-control':'no-cache'}
					response = requests.put(url + '/' + date_filename, file, headers=headers)
					print(url + '/' + date_filename)
					print(response.text)
				
				
					
		return "ok"
		
def upload_data(file_path):
	twitter_obj = Twitter()
	data = list()
	f = open(file_path)
	lines = f.readlines()
	lines = list(map(lambda s: s.strip(),lines))
	for i in lines:
		if i =="":
			continue
		else:
			for a in twitter_obj.nouns(i):
				data.append(a)
			r=re.sub(regex,"",i)
			tokens = nltk.word_tokenize(r)
			tagged = nltk.pos_tag(tokens)
			for item in tagged:
				#print(item)
				if item[1][0]=='N':
					data.append(item[0].lower())
								
	return data

def read_data(folder_path):
	
	data=list() # 폴더마다 명사로 어근 추출 리스트를 가지고 있음 data[0]에는 첫번째 폴더 모든 파일에 대한 어근이 리스트로 있음
	twitter_obj = Twitter()
	imsi=list()
	content = list()
	count=0
	for root,dirs, files in os.walk(folder_path):
		if root != folder_path :
			for fname in files:
				full_fname = os.path.join(root,fname)
				ext = fname.split('.')
				if ext[-1]=='txt':
					f = open(full_fname)
					content = f.read()
				
				elif ext[-1]=='docx':
					content=docx2txt.process(full_fname)
					
				else:
					
					continue

				

			
				for parsed_list in twitter_obj.nouns(content) :	
						
					file_noun_list.append(parsed_list)
					imsi.append(parsed_list)

					
				r=re.sub(regex,"",content)
				tokens = nltk.word_tokenize(r)
				tagged = nltk.pos_tag(tokens)
				#print(tagged)
				for item in tagged:
					#print(item)
					if item[1][0]=='N':
						file_noun_list.append(item[0])
						imsi.append(item[0].lower())
			
			if len(imsi) != 0:
				data.append(imsi)
				
				choosen_folder_path[count]=root
				count=count+1
				folder_file_count.append(files)
			
			imsi=list()
			
						
	return data

def frequency(input_vector, trained_vector):
	frequency=dict()
	
	for words in input_vector: # dict 초기화 
		frequency[words] = 0
	#print(frequency)
	for word in trained_vector:
		if word in input_vector:
			
			frequency[word] += 1
				#print(word)			
	return frequency
#print(data[0])


# v는 유일 단어 갯수 , trained_vector 단일 
# n 은 폴더 내 파일 갯수 
def Log_laplace_probability(input_vector,trained_vector,frequency,v,n,file_num):	 
	denominator = v + len(trained_vector)
	probability = 0
	for word in input_vector:
		probability = probability + math.log((frequency[word]+1)/(denominator))
	print("n/file_num")
	print(n/file_num)
	probability = probability + math.log(n/file_num)
	print("n:")
	print(n)
	print("file_num:")
	print(file_num)

	return probability

#클러스터 초기화 요청.
def requestInitCluster(userId, userToken, folderPath, numberOfCluster):
	print(userId)
	print(userToken)
	print(folderPath)
	print(numberOfCluster)


	
	return "OK"
@app.route('/requestinitcluster', methods = ['POST'])
def requestinitcluster():
	if request.method == 'POST':
		
		uploaded_files = list(request.files.values())
		print("request.form['currentUserId']")
		print(request.form["currentUserId"])
		currentUserId = request.form["currentUserId"]
		currentUserToken = request.form["currentUserToken"]
		currentNumberOfCluster = request.form["currentNumberOfCluster"]
		path_dir = '/home/gxicxigouxa/myproject/users/' + currentUserId + '/textcompare'
		print("UserId: " + currentUserId)
		print("NumberOfCluster: " + currentNumberOfCluster)
		print("request.files: ")
		print(request.files)
		print("uploaded_files: ")
		print (uploaded_files)



		filenames=[]
		dir_list=[]		
		english_stemmer = nltk.stem.SnowballStemmer('english')
		class StemmedTfidfVectorizer(TfidfVectorizer):
			def build_analyzer(self):
				analyzer = super(TfidfVectorizer,self).build_analyzer()
				return lambda doc:(english_stemmer.stem(w) for w in analyzer(doc))
				
		vectorizer = StemmedTfidfVectorizer(min_df=1,max_df=0.5,stop_words='english',decode_error='ignore')			
		total_content=[]
		total_filename=[]
		for file in uploaded_files:		
			if file and allowed_file(file.filename):
				

				filename = file.filename
				filenames.append(filename)
				file.save(os.path.join(path_dir, filename))
				upload_contents=""

				ext = filename.split('.')
				
				if ext[-1]=='txt':
					f= open(path_dir+'/'+filename,'r')
					upload_contents=f.read()

					f.close()
				elif ext[-1]=='docx':
					docx_content = docx2txt.process(path_dir+'/'+filename)
					upload_contents=docx_content
				elif  ext[-1]=='pptx':
					prs =Presentation(path_dir+'/'+filename)
					for slide in prs.slides:
						for shape in slide.shapes:
							if not shape.has_text_frame:
								continue
							for paragraph in shape.text_frame.paragraphs:
								for run in paragraph.runs:
									upload_contents+=run.text+ ' '
				
				else :
					return "no available file extension"+ '  '+str(ext[-1])

				total_content.append(upload_contents);
				total_filename.append(filename)
			else:
				return "not allowed ext"

		print(total_content)
		posts =total_content
		range_n_cluster=[]		
		train_data=posts
		vectorized=vectorizer.fit_transform(train_data)
		num_samples,num_features = vectorized.shape
		print ("ok")
		if len(posts)<=3:
			num_clusters = len(posts)
		else:	
			print("else")
			for i in range(len(posts)-2):
				range_n_cluster.append(i+2);
			print(range_n_cluster)	
			best_k=1;
			best_sil=0;
			for num_clusters in range_n_cluster:
				clusters  = KMeans(n_clusters=num_clusters, n_init=1, verbose=1, random_state=3)
				cluster_labels=cluster.fit_predict(vectorized)
				silhouette_avg = silhouette_score(vectorized,cluster_labels)
				print("n_clustsers =",num_clusters, "The avg silscore : ", silhouette_avg)
				if best_sil<silhouette_avg:
					best_k = num_clusters
					best_sil=silhouette_avg

			num_clusters = best_k
		print(num_clusters)		
		print(total_filename)
		
		km  = KMeans(n_clusters=num_clusters, n_init=1, verbose=1, random_state=3)
		clusterd = km.fit(vectorized)
		print(km.labels_)
		label_set = set(km.labels_)
		for i in range(len(label_set)):
			if not os.path.exists("/home/gxicxigouxa/myproject/users/" + currentUserId+"/textcompare/classified["+str(i)+"]"):	
				os.mkdir("/home/gxicxigouxa/myproject/users/" + currentUserId+"/textcompare/classified["+str(i)+"]")
				os.mkdir("/home/gxicxigouxa/myproject/users/" + currentUserId+"/textcompare/classified["+str(i)+"]/presentation/")
				print("/home/gxicxigouxa/myproject/users/" + currentUserId+"/textcompare/classified"+str(i) + " success.")
				
				requestCreateFolder(currentUserId, currentUserToken, "textcompare", "classified["+str(i)+"]")
				requestCreateFolder(currentUserId, currentUserToken, "textcompare", "classified["+str(i)+"]/presentation")
		for i in range(len(km.labels_)):#ᅟex> [0112201~~]
			for j in range(len(label_set)):#labelset ex> [0,1,2]
				if(km.labels_[i]==j):
					ext = total_filename[i].split(".")
					if ext[-1]=="pptx":
						requestFileUpload(currentUserId, currentUserToken, 'textcompare/classified['+str(i)+']/presentation', uploaded_files[i])
						shutil.move(path_dir+'/'+total_filename[i],"/home/gxicxigouxa/myproject/users/" + currentUserId+"/textcompare/classified["+str(j)+"]/presentation/"+total_filename[i])
					else:
						requestFileUpload(currentUserId, currentUserToken, 'textcompare/classified['+str(i)+']', uploaded_files[i])
						shutil.move(path_dir+'/'+total_filename[i],"/home/gxicxigouxa/myproject/users/" + currentUserId+"/textcompare/classified["+str(j)+"]/"+total_filename[i])	
						
		
		return "OK"
	
@app.route('/textcompare',methods=['POST'])#post with javascript code!!!
#파일을 받아 각 폴더 내부에 있는 파일의 유사도를 분석하여 해당 파일과 가장 알맞는 폴더를 찾아 이동시킨다.
#분류 대상은 텍스트(.txt), MS word(.docx), MS powerpoint(.pptx)이며, 이외에는 분류하지 않고 현재 디렉토리에 저장.
def textcompare():
	if request.method =='POST':
		print("ok1")
		#request로 ID 받아서 pathdir을 /home/gxicxigouxa/user/(숨긴 사용자 ID)/textcompare로 바꾼다.
		#그리고 오픈스택에도 분류된 상태로 올린다.
		
		dir_list=[]
		filenames= []
		filename=''
		upload_contents=''
		print(request)
		'''
		uploaded_files = request.files.getlist('file[]')
		currentUserId = request.form["userId"]
		print(currentUserId)
		'''

		uploaded_files = list(request.files.values())
		print("request.form['currentUserId']")
		print(request.form["currentUserId"])
		currentUserId = request.form["currentUserId"]
		currentUserToken = request.form["currentUserToken"]
		print("UserId: " + currentUserId)
		print("request.files: ")
		print(request.files)
		print("uploaded_files: ")
		print (uploaded_files)
		#path_dir = TXT_DIRECTORY
		path_dir = '/home/gxicxigouxa/myproject/users/' + currentUserId + '/textcompare'
		
		print("path_dir: " + path_dir)
		
		for file in uploaded_files:		
			if file and allowed_file(file.filename):
				

				filename = file.filename
				filenames.append(filename)
				file.save(os.path.join(path_dir, filename))
		

				ext = filename.split('.')
				
				if ext[-1]=='txt':
					f= open(path_dir+'/'+filename,'r')
					upload_contents=f.read()

					f.close()
				elif ext[-1]=='docx':
					docx_content = docx2txt.process(path_dir+'/'+filename)
					upload_contents=docx_content
				elif  ext[-1]=='pptx':
					prs =Presentation(path_dir+'/'+filename)
					for slide in prs.slides:
						for shape in slide.shapes:
							if not shape.has_text_frame:
								continue
							for paragraph in shape.text_frame.paragraphs:
								for run in paragraph.runs:
									upload_contents+=run.text+ ' '
				else :
					return "no available file extension"+ '  '+str(ext[-1])

			else:
				return "not allowed ext"
			
			###--------------------------parsing only directory (not file)-----------------------------------------
			for file in os.listdir(path_dir):
				if os.path.isdir(path_dir+'/'+file):
					dir_list.append(path_dir+'/'+file)

			dir_list.sort()
			###--------------------------parsing only directory (not file)-------------------------------------------
			
			
			total_score=[]
			total_content=[]
			#same index with dir_list
			


			#vectorizer = CountVectorizer(min_df=1)
			vectorizer=StemmedCountVectorizer(min_df=1,max_df=0.9, stop_words=mystopword)
			#parsing the all content in the test file not in the presentation file  pptxfile is in the another paht (ex > not /test1   /test1/presentation)
			for i in range(len(dir_list)):
				for file in os.listdir(dir_list[i]):
				
					if os.path.isfile(dir_list[i]+'/'+file):	# if it is file 
						ext =file.split('.')
						if ext[-1]=='txt':
							f= open(dir_list[i]+'/'+file)
							total_content.append(f.read())
						if ext[-1]=='docx':
							total_content.append(docx2txt.process(dir_list[i]+'/'+file))

					else:# if it is directory
						continue			
			#parsing the all content

			                                                
			for i in range(len(total_content)):
				X_train=vectorizer.fit_transform(total_content)
				num_samples,num_features=X_train.shape
				new_post=upload_contents
				new_post_vec = vectorizer.transform([new_post])

			best_doc = None
			best_dist = sys.maxsize
			best_i=None
			for i,post in enumerate(total_content):
				#if post == new_post:
				#	continue

				post_vec = X_train.getrow(i)
				d=dist_norm(post_vec,new_post_vec)
				total_score.append(d)
				if d<best_dist:
					best_dist=d
					best_i=i

			dir_file_num=[]
			# directory's number of files!

			for i in range(len(dir_list)):
				dir_file_num.append(num_only_file(dir_list[i]))#do not count the  directory , just file count(because of presentation)


			sum =0
			#find the file's closest directory location
			for i in range(len(dir_file_num)):
				sum+=dir_file_num[i] 
				if(sum>=best_i+1):
					closest_location=i
					break				
			

			most= dir_list[closest_location]

			
			if(filename.split('.')[-1]=='pptx'):# upload file extension is pptx !!! go to presentation
				shutil.move(path_dir+'/'+filename,most+"/"+"presentation/"+filename)
				print("most: " + most)
				#cuttedMost = most.split("/home/gxicxigouxa/myproject/users/" + currentUserId + "/textcompare/")
				cuttedMost = most.split(path_dir + "/")
				print("cuttedMost: ")
				print(cuttedMost)
				#cuttedMost[1]
				url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ currentUserId + "/textcompare/" + cuttedMost[1] + "/presentation/"
				print(url)
				print("file name: " + filename)
				headers  ={'X-File-Name':filename, 'x-auth-token':currentUserToken,'content-type':'text/html', 'cache-control':'no-cache'}
				response = requests.put(url + '/' + filename, file, headers=headers)
				print(url + '/' + filename)
				print(response.text)
			else:
				shutil.move(path_dir+'/'+filename,most+"/"+filename)
				print("most: " + most)
				cuttedMost = most.split				("/home/gxicxigouxa/myproject/users/" + currentUserId + "/textcompare/")
				cuttedMost = most.split(path_dir + "/")
				print("cuttedMost: ")
				print(cuttedMost)
				#cuttedMost[1]
				url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ currentUserId + "/textcompare/" + cuttedMost[1]
				print(url)
				print("file name: " + filename)
				headers  ={'X-File-Name':filename, 'x-auth-token':currentUserToken,'content-type':'text/html', 'cache-control':'no-cache'}
				response = requests.put(url + '/' + filename, file, headers=headers)
				print(url + '/' + filename)
				print(response.text)
		
			#return "The closest is " +str(closest_location+1)+"'st directory! So the path is '"+str(most)+"'   totalscore " + str(total_score)+"total content : " + str(total_content)		
		#after this code ----------------- front section modified code
	textdir_list=[]
	textdir_list = os.listdir(path_dir)
	return "ok"

@app.route('/dbincomeexpect', methods=['POST', 'GET'])
def dbincomeexpect():
	conn = mysql.connect()
	cursor=conn.cursor()
	cursor.execute("select * from userinfotable")
	
	result =[]
	columns= tuple( [d[0] for d in cursor.description] )

	
	for row in cursor:
		result.append(dict(zip(columns, row)))
	
	numOfJson= len(result)
	numberOfAccount = 0
	id =[]
	usedmonth=[]
	grade=[]
	totalamount=[]
	numofregist=[]
	averagefee=[]
	data= [] 
	target = []
	userid =[]
	xdata = []
	ytarget=[]

	for i in range(numOfJson):
		if result[i]["id"] == "admin":
			print(result[i]["id"])
			continue
		id.append(result[i]["id"])
		usedmonth.append(result[i]["paymentday"])
		grade.append(result[i]["rating"])
		totalamount.append(result[i]["totalamount"])
		numofregist.append(result[i]["enrollmentnumber"])
		averagefee.append(result[i]["yearaverageamount"])
		tmp=str(usedmonth[numberOfAccount])+' '+str(grade[numberOfAccount])+' '+str(totalamount[numberOfAccount])+ ' '+str(numofregist[numberOfAccount])
		data.append(tmp.split(' '))
		target.append(averagefee[numberOfAccount])
		userid.append(id[numberOfAccount])
		print("OK" + str(i))
		numberOfAccount += 1
	# ------------------------------------------------array element's string type to float!!!!!!!!!!!!!!------------------------------------------------
	for j in range(numberOfAccount):
		xdata.append([float(i) for i in data[j]])

	print(len(xdata))
	print(numberOfAccount)
	ytarget =[float(i) for i in target]
	
	# like map (int,list)
	# ------------------------------------------------array element's string type to float!!!!!!!!!!!!!!------------------------------------------------
	
	lr = LinearRegression()
	en = ElasticNet(alpha=0.5,precompute=False)

	
	
	y= np.transpose(np.atleast_2d(ytarget))#transpose ydata to two dimentional array
	

	lr.fit(xdata,y)
	
	lrp =lr.predict(xdata)
	
	

	kf = KFold(len(xdata),n_folds=5)

	enp = np.zeros_like(y)
	###---------------------------------important!!!!!!!!!!!!!!!!!!-----------------------------------------
	enpxdata= np.zeros_like(xdata)###important!!!!!!!!!!!!!!!!


	for i in range(len(xdata)):
		enpxdata[i]=xdata[i]
	#xdata => [[1,2,3],[4.5.6]]  enpxdata=> [[1,2,3] [4,5,6]]           no comma!!!!!	
	###---------------------------------important!!!!!!!!!!!!!!!!!!-----------------------------------------
	for train,test in kf:
		en.fit(enpxdata[train],y[train])
		enpred=np.transpose(np.atleast_2d(en.predict(enpxdata[test])))#it is very important!!!
		enp[test]=enpred
	
	EnLrswap(enp,lrp)#minus element must be swap!

	
	sum = 0
	personalExpectedIncomeList = []
	for i in range(len(enp)):
		sum+=enp[i,0]
		personalExpectedIncomeList.append(round(enp[i,0],2))

	return jsonify({"data":personalExpectedIncomeList,"id":userid})

'''
@app.route('/dbincomeexpect')
def dbincomeexpect():
	conn = mysql.connect()
	cursor=conn.cursor()
	cursor.execute("select * from jhjtable")
	
	result =[]
	columns= tuple( [d[0] for d in cursor.description] )

	
	for row in cursor:
		result.append(dict(zip(columns, row)))
	
	numOfJson= len(result)
	id =[]
	usedmonth=[]
	grade=[]
	totalamount=[]
	numofregist=[]
	averagefee=[]
	data= [] 
	target = []
	userid =[]
	xdata = []
	ytarget=[]

	for i in range(numOfJson):
		id.append(result[i]["id"])
		usedmonth.append(result[i]["usedmonth"])
		grade.append(result[i]["grade"])
		totalamount.append(result[i]["totalamount"])
		numofregist.append(result[i]["numofregist"])
		averagefee.append(result[i]["averagefee"])
		tmp=str(usedmonth[i])+' '+str(grade[i])+' '+str(totalamount[i])+ ' '+str(numofregist[i])
		data.append(tmp.split(' '))
		target.append(averagefee[i])
		userid.append(id[i])
		
	# ------------------------------------------------array element's string type to float!!!!!!!!!!!!!!------------------------------------------------
	for j in range(numOfJson):
		xdata.append([float(i) for i in data[j]])

		
	ytarget =[float(i) for i in target]
	
	# like map (int,list)
	# ------------------------------------------------array element's string type to float!!!!!!!!!!!!!!------------------------------------------------
	
	lr = LinearRegression()
	kf = KFold(len(xdata),n_folds=5)
	#en = ElasticNet(alpha=0.5,precompute=False)


	met = ElasticNetCV(l1_ratio=l1_ratio,n_jobs=-1)
	y= np.transpose(np.atleast_2d(ytarget))#transpose ydata to two dimentional array
	

	lr.fit(xdata,y)
	
	lrp =lr.predict(xdata)
	
	


	

	enp = np.zeros_like(y)
	###---------------------------------important!!!!!!!!!!!!!!!!!!-----------------------------------------
	enpxdata= np.zeros_like(xdata)###important!!!!!!!!!!!!!!!!


	for i in range(len(xdata)):
		enpxdata[i]=xdata[i]
	#xdata => [[1,2,3],[4.5.6]]  enpxdata=> [[1,2,3] [4,5,6]]           no comma!!!!!	
	###---------------------------------important!!!!!!!!!!!!!!!!!!-----------------------------------------
	for train,test in kf:
		met.fit(enpxdata[train],y[train])
		enpred=np.transpose(np.atleast_2d(en.predict(enpxdata[test])))#it is very important!!!
		enp[test]=enpred
	
	EnLrswap(enp,lrp)#minus element must be swap!

	
	sum = 0
	for i in range(len(enp)):
	
		sum+=enp[i,0]
	
			
	return str(sum)
@app.route('/dbincomepersonal')
def dbincomepersonal():
	conn = mysql.connect()
	cursor=conn.cursor()
	cursor.execute("select * from jhjtable")
	
	result =[]
	columns= tuple( [d[0] for d in cursor.description] )

	
	for row in cursor:
		result.append(dict(zip(columns, row)))
	
	numOfJson= len(result)
	id =[]
	usedmonth=[]
	grade=[]
	totalamount=[]
	numofregist=[]
	averagefee=[]
	data= [] 
	target = []
	userid =[]
	xdata = []
	ytarget=[]

	for i in range(numOfJson):
		id.append(result[i]["id"])
		usedmonth.append(result[i]["usedmonth"])
		grade.append(result[i]["grade"])
		totalamount.append(result[i]["totalamount"])
		numofregist.append(result[i]["numofregist"])
		averagefee.append(result[i]["averagefee"])
		tmp=str(usedmonth[i])+' '+str(grade[i])+' '+str(totalamount[i])+ ' '+str(numofregist[i])
		data.append(tmp.split(' '))
		target.append(averagefee[i])
		userid.append(id[i])
		
	# ------------------------------------------------array element's string type to float!!!!!!!!!!!!!!------------------------------------------------
	for j in range(numOfJson):
		xdata.append([float(i) for i in data[j]])

		
	ytarget =[float(i) for i in target]
	
	# like map (int,list)
	# ------------------------------------------------array element's string type to float!!!!!!!!!!!!!!------------------------------------------------
	
	lr = LinearRegression()
	en = ElasticNet(alpha=0.5,precompute=False)

	
	
	y= np.transpose(np.atleast_2d(ytarget))#transpose ydata to two dimentional array
	

	lr.fit(xdata,y)
	
	lrp =lr.predict(xdata)
	
	

	kf = KFold(len(xdata),n_folds=5)

	enp = np.zeros_like(y)
	###---------------------------------important!!!!!!!!!!!!!!!!!!-----------------------------------------
	enpxdata= np.zeros_like(xdata)###important!!!!!!!!!!!!!!!!


	for i in range(len(xdata)):
		enpxdata[i]=xdata[i]
	#xdata => [[1,2,3],[4.5.6]]  enpxdata=> [[1,2,3] [4,5,6]]           !!!!!<no comma>!!!!!
	#if I do not this  knn expect didn't operate
	###---------------------------------important!!!!!!!!!!!!!!!!!!-----------------------------------------
	for train,test in kf:
		en.fit(enpxdata[train],y[train])
		enpred=np.transpose(np.atleast_2d(en.predict(enpxdata[test])))#it is very important!!!
		enp[test]=enpred
	
	EnLrswap(enp,lrp)#minus element must be swap!

	
	showall= ''
	enpshowall = ""
	i=0
	for i in range(len(data)):
		showall+=str(enp[i,0])+'\n'
		
		
	return str(showall)
'''

		#
@app.route('/membertable',methods=['POST','GET'])
def	membertable():
	if request.method =='GET':
		conn =mysql.connect()
		cursor =conn.cursor()
		query = "select * from userinfotable;"
		cursor.execute(query)
		conn.commit()
			
		result =[]
		columns= tuple( [d[0] for d in cursor.description] )
			
		for row in cursor:
			result.append(dict(zip(columns, row)))
			print (str(result))
		numOfJson= len(result)
		admin_excluded_mem_num = 0
		id =[]
		pwd=[]
		birth=[]
		email=[]
		usedmonth=[]
		grade=[]
		totalamount=[]
		numofregist=[]
		averagefee=[]
		data= [] 
		expiretime=[]

		for i in range(numOfJson):
			if result[i]["id"] == "admin":
				print(result[i]["id"])
				continue
			id.append(result[i]["id"])
			pwd.append(result[i]["pwd"])
			birth.append(result[i]["birth"])
			email.append(result[i]["email"])
			usedmonth.append(result[i]["paymentday"])
			expiretime.append(result[i]["expiretime"])
			grade.append(result[i]["rating"])
			totalamount.append(result[i]["totalamount"])
			numofregist.append(result[i]["enrollmentnumber"])
			averagefee.append(result[i]["yearaverageamount"])

			admin_excluded_mem_num += 1

		return jsonify({"id":id,"pwd":pwd,"birth":birth,"email":email,"usedmonth":usedmonth,"expiretime":expiretime,"grade":grade,"totalamount":totalamount,"numofregist":numofregist,"averagefee":averagefee})

@app.route('/textcomparefilemove',methods=['POST','GET'])
def textcomparefilemove():
	ALL_PATH=[]
	#토큰에 대한 컨테이너 목록 요청
	print(session["userId"])
	container_url ='http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ session["userId"]
	headers  ={'x-auth-token':session["token"],'content-type':'application/json'}
	response = requests.get(container_url,headers=headers)
	print(container_url)
	containerlist = response.text.split("\n")

	for container in containerlist:
		if container !="textcompare":
			continue
		filelist_url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ session["userId"]+'/'+container
		response=requests.get(filelist_url,headers=headers)
		if response.text=="":
			ALL_PATH.append(container+'/')# because When container is empty there is blank response  so it cannot go into the sencond depth  for loop append the path here
		path_list=response.text.split("\n")

		for path in path_list:
			if path =='':
				continue
			folder_file_list=path.split('/')# when split, folder's last index is '' and file is ~~.ext
			if folder_file_list[-1] == '':# last index is represent where it is file or folder. append only folder 
				ALL_PATH.append(container+'/'+path)

	
	#return str(ALL_PATH)
	return jsonify({"pathList":ALL_PATH})
english_stemmer = nltk.stem.SnowballStemmer('english')
class StemmedTfidfVectorizer(TfidfVectorizer):
	def build_analyzer(self):
		analyzer = super(TfidfVectorizer,self).build_analyzer()
		return lambda doc:(english_stemmer.stem(w) for w in analyzer(doc))
	
@app.route('/knn',methods =['POST'])
def makingtrainingdata():
		
	if request.method =='POST':
		vectorizer = StemmedTfidfVectorizer(min_df=2,max_df=0.5,stop_words='english',decode_error='ignore')
		upload_contents=''
		uploaded_files = list(request.files.values())
		print("request.form['currentUserId']")
		print(request.form["currentUserId"])
		currentUserId = request.form["currentUserId"]
		currentUserToken = request.form["currentUserToken"]
		print("UserId: " + currentUserId)
		print("request.files: ")
		print(request.files)
		print("uploaded_files: ")
		print (uploaded_files)
		path_dir = '/home/gxicxigouxa/myproject/users/' + currentUserId + '/textcompare'
		upload_contents_list=list()
		for file in uploaded_files:		
			if file and allowed_file(file.filename):
				

				filename = file.filename
			
				file.save(os.path.join(path_dir, filename))
		

				ext = filename.split('.')
				
				if ext[-1]=='txt':
					f= open(path_dir+'/'+filename,'r')
					upload_contents=f.read()
					
					f.close()
				elif ext[-1]=='docx':
					docx_content = docx2txt.process(path_dir+'/'+filename)
					upload_contents=docx_content
				elif  ext[-1]=='pptx':
					prs =Presentation(path_dir+'/'+filename)
					for slide in prs.slides:
						for shape in slide.shapes:
							if not shape.has_text_frame:
								continue
							for paragraph in shape.text_frame.paragraphs:
								for run in paragraph.runs:
									upload_contents+=run.text+ ' '
				
			
				else :
					return "no available file extension"+ '  '+str(ext[-1])
				upload_contents_list.append(upload_contents)
				train_data = upload_contents_list
				vectorized=vectorizer.fit_train(train_data)
				

		print(str(upload_contents_list))


		return "ok"

@app.route('/storagefilemove',methods=['POST','GET'])
def storagefilemove():
	ALL_PATH=[]
	#토큰에 대한 컨테이너 목록 요청
	print(session["userId"])
	container_url ='http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ session["userId"]
	headers  ={'x-auth-token':session["token"],'content-type':'application/json'}
	response = requests.get(container_url,headers=headers)
	print(container_url)
	containerlist = response.text.split("\n")

	for container in containerlist:
		if (container =="textcompare" or container =="malware" or container == ""):
			continue
		filelist_url = 'http://' + OPENSTACK_IP + ':8080/v1/AUTH_'+ session["userId"]+'/'+container
		response=requests.get(filelist_url,headers=headers)
		if response.text=="":
			ALL_PATH.append(container+'/')# because When container is empty there is blank response  so it cannot go into the sencond depth  for loop append the path here
		path_list=response.text.split("\n")

		for path in path_list:
			if path =='':
				continue
			folder_file_list=path.split('/')# when split, folder's last index is '' and file is ~~.ext
			if folder_file_list[-1] == '':# last index is represent where it is file or folder. append only folder 
				ALL_PATH.append(container+'/'+path)

	
	return jsonify({"pathList":ALL_PATH})
	#return str(ALL_PATH)







if __name__ =='__main__':
   app.run(host='0.0.0.0',port=9999)	